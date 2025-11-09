#!/usr/bin/env python3
"""
Create YTL Cement PowerPoint with Native Shapes (Not Images)
All diagrams are built using PowerPoint rectangles, text boxes, and connectors
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_shape_box(slide, left, top, width, height, text, fill_color, text_size=12, bold=False):
    """Add a rectangular shape with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    # Set fill color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*fill_color)

    # Set border
    line = shape.line
    line.color.rgb = RGBColor(100, 100, 100)
    line.width = Pt(1)

    # Add text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Format text
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.name = 'Arial'

    return shape

def add_connector(slide, shape1, shape2, label=""):
    """Add a connector between two shapes"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        shape1.left + shape1.width // 2,
        shape1.top + shape1.height,
        shape2.left + shape2.width // 2,
        shape2.top
    )
    connector.line.color.rgb = RGBColor(100, 100, 100)
    connector.line.width = Pt(2)
    return connector

def add_note_box(slide, left, top, width, height, title, bullets):
    """Add a note/annotation box"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    # Light yellow background
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 220)

    # Border
    line = shape.line
    line.color.rgb = RGBColor(200, 200, 100)
    line.width = Pt(1)

    # Add text
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    # Title
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(10)

    # Bullets
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(9)

    return shape

# ===== SLIDE 1: Title =====
add_title_slide(
    prs,
    "YTL Cement - SAP S/4HANA Transformation",
    "Business Case & Technical Architecture\nFully Editable Native PowerPoint Diagrams"
)

# ===== SLIDE 2: Capability Map AS-IS =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = add_shape_box(
    slide, Inches(0.5), Inches(0.3), Inches(15), Inches(0.6),
    "YTL Cement - Business Capability Map (Current State - AS-IS)",
    (255, 255, 255), text_size=24, bold=True
)

# Main container
container = add_shape_box(
    slide, Inches(0.5), Inches(1), Inches(9), Inches(7),
    "", (245, 245, 245), text_size=10
)

# Strategic Level
strategic_label = add_shape_box(
    slide, Inches(0.7), Inches(1.2), Inches(8.6), Inches(0.4),
    "STRATEGIC LEVEL", (240, 240, 240), text_size=11, bold=True
)

# Strategic boxes
p2p = add_shape_box(slide, Inches(0.9), Inches(1.8), Inches(2), Inches(0.8),
    "Procure-to-Pay", (255, 244, 230), text_size=11)
o2c = add_shape_box(slide, Inches(3.1), Inches(1.8), Inches(2), Inches(0.8),
    "Order-to-Cash", (255, 244, 230), text_size=11)
mts = add_shape_box(slide, Inches(5.3), Inches(1.8), Inches(2), Inches(0.8),
    "Make-to-Stock", (255, 244, 230), text_size=11)
fin = add_shape_box(slide, Inches(7.5), Inches(1.8), Inches(1.7), Inches(0.8),
    "Finance &\nReporting", (255, 244, 230), text_size=11)

# Business Capability Level
cap_label = add_shape_box(
    slide, Inches(0.7), Inches(2.9), Inches(8.6), Inches(0.4),
    "BUSINESS CAPABILITY LEVEL", (240, 240, 240), text_size=11, bold=True
)

# Capability boxes (RED - Manual, Error-Prone)
invoice_cap = add_shape_box(slide, Inches(0.9), Inches(3.5), Inches(1.6), Inches(0.9),
    "Invoice\nProcessing", (255, 107, 107), text_size=10, bold=True)
gl_cap = add_shape_box(slide, Inches(2.7), Inches(3.5), Inches(1.6), Inches(0.9),
    "GL\nReconciliation", (255, 107, 107), text_size=10, bold=True)
vendor_cap = add_shape_box(slide, Inches(4.5), Inches(3.5), Inches(1.6), Inches(0.9),
    "Vendor\nManagement", (255, 107, 107), text_size=10, bold=True)
sod_cap = add_shape_box(slide, Inches(6.3), Inches(3.5), Inches(1.6), Inches(0.9),
    "SoD\nCompliance", (255, 107, 107), text_size=10, bold=True)
einvoice_cap = add_shape_box(slide, Inches(8.1), Inches(3.5), Inches(1.1), Inches(0.9),
    "e-Invoice\nCompliance", (176, 176, 176), text_size=9, bold=True)

# Functional Level
func_label = add_shape_box(
    slide, Inches(0.7), Inches(4.7), Inches(8.6), Inches(0.4),
    "FUNCTIONAL LEVEL", (240, 240, 240), text_size=11, bold=True
)

# Functional boxes (Light RED)
po_func = add_shape_box(slide, Inches(0.9), Inches(5.3), Inches(1.5), Inches(0.8),
    "Create PO\nCreate PR", (255, 179, 179), text_size=9)
match_func = add_shape_box(slide, Inches(2.6), Inches(5.3), Inches(1.5), Inches(0.8),
    "Match Invoice\nto PO/GR", (255, 179, 179), text_size=9)
gl_func = add_shape_box(slide, Inches(4.3), Inches(5.3), Inches(1.5), Inches(0.8),
    "Post GL\nEntries", (255, 179, 179), text_size=9)
sod_func = add_shape_box(slide, Inches(6.0), Inches(5.3), Inches(1.5), Inches(0.8),
    "Detect SoD\nViolations", (255, 179, 179), text_size=9)
xml_func = add_shape_box(slide, Inches(7.7), Inches(5.3), Inches(1.5), Inches(0.8),
    "Manual XML\nGeneration", (208, 208, 208), text_size=9)

# Add notes on the right
note1 = add_note_box(slide, Inches(10), Inches(1.5), Inches(5.5), Inches(1.2),
    "🔴 RED = Manual, Error-Prone",
    [
        "Current: 28K invoices/month",
        "Issues: 1-2% undetected mismatches",
        "Cost: $280K-560K annually"
    ])

note2 = add_note_box(slide, Inches(10), Inches(2.9), Inches(5.5), Inches(1.2),
    "🔴 RED = Manual Process",
    [
        "Current: 8-day close cycle",
        "Issues: Manual variance investigation",
        "Cost: $192K annually (labor)"
    ])

note3 = add_note_box(slide, Inches(10), Inches(4.3), Inches(5.5), Inches(1.2),
    "🔴 RED = Excel-Based",
    [
        "Current: 8,000 vendors",
        "Issues: 23% duplicates",
        "Cost: $200K-400K loss annually"
    ])

note4 = add_note_box(slide, Inches(10), Inches(5.7), Inches(5.5), Inches(1.1),
    "🟡 YELLOW = Partial",
    [
        "Current: Manual submission",
        "Issues: Error-prone, not scalable",
        "Deadline: Q2 2026 (LHDN MyInvois)"
    ])

# Legend
legend = add_shape_box(slide, Inches(10), Inches(7), Inches(5.5), Inches(0.9),
    "", (255, 255, 255), text_size=9)
legend_text = legend.text_frame
legend_text.clear()
p1 = legend_text.paragraphs[0]
p1.text = "Legend:"
p1.font.bold = True
p1.font.size = Pt(10)

legend_items = [
    ("🔴 RED", "Manual, Error-Prone, High Risk"),
    ("🟡 YELLOW", "Partially Automated"),
    ("🟢 GREEN", "Fully Automated, Optimized")
]
for emoji, desc in legend_items:
    p = legend_text.add_paragraph()
    p.text = f"{emoji} = {desc}"
    p.font.size = Pt(8)


# ===== SLIDE 3: Capability Map TO-BE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = add_shape_box(
    slide, Inches(0.5), Inches(0.3), Inches(15), Inches(0.6),
    "YTL Cement - Business Capability Map (Target State - TO-BE)",
    (255, 255, 255), text_size=24, bold=True
)

# Main container
container = add_shape_box(
    slide, Inches(0.5), Inches(1), Inches(9), Inches(7),
    "", (245, 245, 245), text_size=10
)

# Strategic Level
strategic_label = add_shape_box(
    slide, Inches(0.7), Inches(1.2), Inches(8.6), Inches(0.4),
    "STRATEGIC LEVEL", (240, 240, 240), text_size=11, bold=True
)

# Strategic boxes (same as AS-IS)
p2p = add_shape_box(slide, Inches(0.9), Inches(1.8), Inches(2), Inches(0.8),
    "Procure-to-Pay", (255, 244, 230), text_size=11)
o2c = add_shape_box(slide, Inches(3.1), Inches(1.8), Inches(2), Inches(0.8),
    "Order-to-Cash", (255, 244, 230), text_size=11)
mts = add_shape_box(slide, Inches(5.3), Inches(1.8), Inches(2), Inches(0.8),
    "Make-to-Stock", (255, 244, 230), text_size=11)
fin = add_shape_box(slide, Inches(7.5), Inches(1.8), Inches(1.7), Inches(0.8),
    "Finance &\nReporting", (255, 244, 230), text_size=11)

# Business Capability Level
cap_label = add_shape_box(
    slide, Inches(0.7), Inches(2.9), Inches(8.6), Inches(0.4),
    "BUSINESS CAPABILITY LEVEL", (240, 240, 240), text_size=11, bold=True
)

# Capability boxes (GREEN - Automated)
invoice_cap = add_shape_box(slide, Inches(0.9), Inches(3.5), Inches(1.6), Inches(0.9),
    "Invoice\nProcessing", (102, 187, 106), text_size=10, bold=True)
gl_cap = add_shape_box(slide, Inches(2.7), Inches(3.5), Inches(1.6), Inches(0.9),
    "GL\nReconciliation", (102, 187, 106), text_size=10, bold=True)
vendor_cap = add_shape_box(slide, Inches(4.5), Inches(3.5), Inches(1.6), Inches(0.9),
    "Vendor\nManagement", (102, 187, 106), text_size=10, bold=True)
sod_cap = add_shape_box(slide, Inches(6.3), Inches(3.5), Inches(1.6), Inches(0.9),
    "SoD\nCompliance", (102, 187, 106), text_size=10, bold=True)
einvoice_cap = add_shape_box(slide, Inches(8.1), Inches(3.5), Inches(1.1), Inches(0.9),
    "e-Invoice\nCompliance", (102, 187, 106), text_size=9, bold=True)

# Functional Level
func_label = add_shape_box(
    slide, Inches(0.7), Inches(4.7), Inches(8.6), Inches(0.4),
    "FUNCTIONAL LEVEL", (240, 240, 240), text_size=11, bold=True
)

# Functional boxes (Light GREEN)
po_func = add_shape_box(slide, Inches(0.9), Inches(5.3), Inches(1.5), Inches(0.8),
    "Automated PO\nGeneration", (168, 213, 186), text_size=9)
match_func = add_shape_box(slide, Inches(2.6), Inches(5.3), Inches(1.5), Inches(0.8),
    "Continuous\n3-Way Matching", (168, 213, 186), text_size=9)
gl_func = add_shape_box(slide, Inches(4.3), Inches(5.3), Inches(1.5), Inches(0.8),
    "Real-Time GL\nPosting", (168, 213, 186), text_size=9)
sod_func = add_shape_box(slide, Inches(6.0), Inches(5.3), Inches(1.5), Inches(0.8),
    "Real-Time SoD\nMonitoring", (168, 213, 186), text_size=9)
xml_func = add_shape_box(slide, Inches(7.7), Inches(5.3), Inches(1.5), Inches(0.8),
    "Automated\nMyInvois", (168, 213, 186), text_size=9)

# MVP Module Layer
mvp_label = add_shape_box(
    slide, Inches(0.7), Inches(6.4), Inches(8.6), Inches(0.4),
    "MVP MODULE LAYER", (227, 242, 253), text_size=11, bold=True
)

# MVP boxes (Blue)
mvp_invoice = add_shape_box(slide, Inches(0.9), Inches(7.0), Inches(1.5), Inches(0.7),
    "Invoice\nMatching", (66, 165, 245), text_size=9, bold=True)
mvp_gl = add_shape_box(slide, Inches(2.6), Inches(7.0), Inches(1.5), Inches(0.7),
    "GL Anomaly\nDetector", (66, 165, 245), text_size=9, bold=True)
mvp_vendor = add_shape_box(slide, Inches(4.3), Inches(7.0), Inches(1.5), Inches(0.7),
    "Vendor Data\nQuality", (66, 165, 245), text_size=9, bold=True)
mvp_sod = add_shape_box(slide, Inches(6.0), Inches(7.0), Inches(1.5), Inches(0.7),
    "SoD\nAnalyzer", (66, 165, 245), text_size=9, bold=True)
mvp_einvoice = add_shape_box(slide, Inches(7.7), Inches(7.0), Inches(1.5), Inches(0.7),
    "e-Invoice\nModule", (66, 165, 245), text_size=9, bold=True)

# Add notes on the right
note1 = add_note_box(slide, Inches(10), Inches(1.5), Inches(5.5), Inches(1.2),
    "🟢 GREEN = Automated",
    [
        "Target: 99.7% auto-match rate",
        "Benefit: $280K-560K savings",
        "Timeline: Continuous monitoring"
    ])

note2 = add_note_box(slide, Inches(10), Inches(2.9), Inches(5.5), Inches(1.2),
    "🟢 GREEN = Real-Time",
    [
        "Target: 3-4 day close",
        "Benefit: $192K labor savings",
        "Timeline: Real-time anomaly detection"
    ])

note3 = add_note_box(slide, Inches(10), Inches(4.3), Inches(5.5), Inches(1.2),
    "🟢 GREEN = Data Quality",
    [
        "Target: <5% duplicates",
        "Benefit: $200K-400K recovery",
        "Timeline: Automated de-duplication"
    ])

note4 = add_note_box(slide, Inches(10), Inches(5.7), Inches(5.5), Inches(1.1),
    "🟢 GREEN = Automated",
    [
        "Target: 100% LHDN compliance",
        "Benefit: 0.5 FTE labor savings",
        "Timeline: Automated submission"
    ])


# ===== SLIDE 4: Gap Analysis =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = add_shape_box(
    slide, Inches(0.5), Inches(0.3), Inches(15), Inches(0.6),
    "Gap Analysis: Current vs Target State",
    (255, 255, 255), text_size=24, bold=True
)

gaps = [
    {
        "name": "Invoice Processing",
        "y": 1.2,
        "as_is": ["Manual 3-Way Match", "28K invoices/month", "1-2% errors", "$280K-560K/yr cost"],
        "solution": ["Invoice Matching MVP", "Continuous monitoring", "99.7% accuracy", "Pattern detection"],
        "to_be": ["Automated Match", "ROI: <1 month", "Effort: 4-5 weeks", "Value: $280K-560K"]
    },
    {
        "name": "GL Reconciliation",
        "y": 2.7,
        "as_is": ["8-day close cycle", "Manual variance", "2-3 FTE per close", "$192K/yr cost"],
        "solution": ["GL Anomaly Detector", "Real-time monitoring", "4-day close target", "Auto-detection"],
        "to_be": ["Real-Time Recon", "ROI: 1-2 months", "Effort: 5-6 weeks", "Value: $192K/yr"]
    },
    {
        "name": "Vendor Data Quality",
        "y": 4.2,
        "as_is": ["8,000 vendors", "23% duplicates", "Manual management", "$200K-400K/yr"],
        "solution": ["Vendor Quality MVP", "Auto de-duplication", "Quality scoring", "Compliance checks"],
        "to_be": ["Automated Master", "ROI: 1-2 months", "Effort: 6-8 weeks", "Value: $200K-400K"]
    },
    {
        "name": "SoD Compliance",
        "y": 5.7,
        "as_is": ["Quarterly audits", "6-month lag", "Audit findings", "HIGH risk"],
        "solution": ["SoD Analyzer MVP", "Real-time monitor", "25 key rules", "Auto-escalation"],
        "to_be": ["Real-Time SoD", "ROI: 2-3 months", "Effort: 5-6 weeks", "Value: Risk mitigation"]
    },
]

for gap in gaps:
    # Gap header
    gap_header = add_shape_box(
        slide, Inches(0.5), Inches(gap["y"]), Inches(15), Inches(0.4),
        f"GAP: {gap['name']}",
        (200, 200, 200), text_size=12, bold=True
    )

    # AS-IS box (RED)
    asis_box = add_shape_box(
        slide, Inches(0.5), Inches(gap["y"] + 0.5), Inches(4.5), Inches(1),
        "",
        (255, 153, 153), text_size=9
    )
    asis_text = asis_box.text_frame
    asis_text.clear()
    p = asis_text.paragraphs[0]
    p.text = "AS-IS"
    p.font.bold = True
    p.font.size = Pt(10)
    for item in gap["as_is"]:
        p = asis_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)

    # SOLUTION box (YELLOW)
    sol_box = add_shape_box(
        slide, Inches(5.5), Inches(gap["y"] + 0.5), Inches(4.5), Inches(1),
        "",
        (255, 255, 153), text_size=9
    )
    sol_text = sol_box.text_frame
    sol_text.clear()
    p = sol_text.paragraphs[0]
    p.text = "SOLUTION"
    p.font.bold = True
    p.font.size = Pt(10)
    for item in gap["solution"]:
        p = sol_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)

    # TO-BE box (GREEN)
    tobe_box = add_shape_box(
        slide, Inches(10.5), Inches(gap["y"] + 0.5), Inches(5), Inches(1),
        "",
        (153, 255, 153), text_size=9
    )
    tobe_text = tobe_box.text_frame
    tobe_text.clear()
    p = tobe_text.paragraphs[0]
    p.text = "TO-BE"
    p.font.bold = True
    p.font.size = Pt(10)
    for item in gap["to_be"]:
        p = tobe_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(8)

# Summary note at bottom
summary = add_note_box(slide, Inches(0.5), Inches(7.2), Inches(15), Inches(1.2),
    "Total Gap Closure Summary:",
    [
        "5 Major Gaps → 5 MVP Modules",
        "Total Investment: $160K-240K (MVP modules)",
        "Total Year 1 Benefit: $722K-1.2M",
        "Payback Period: 10-16 months",
        "Implementation: 18-week Phase 1"
    ])


# ===== SLIDE 5: ROI Summary =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = add_shape_box(
    slide, Inches(0.5), Inches(0.3), Inches(15), Inches(0.6),
    "ROI & Annual Business Value",
    (255, 255, 255), text_size=24, bold=True
)

# Investment section (LEFT)
inv_label = add_shape_box(
    slide, Inches(0.5), Inches(1.2), Inches(7), Inches(0.5),
    "INVESTMENT REQUIRED", (255, 230, 230), text_size=14, bold=True
)

investments = [
    ("S/4HANA License (1 year)", "$200K-250K"),
    ("Infrastructure/Hardware", "$150K-200K"),
    ("Implementation Services (1,000 MD)", "$200K-300K"),
    ("MVP Modules (5 modules)", "$160K-240K"),
    ("Change Management & Training", "$100K-150K"),
]

y_pos = 1.8
for item, cost in investments:
    inv_box = add_shape_box(
        slide, Inches(0.5), Inches(y_pos), Inches(5), Inches(0.5),
        item, (255, 240, 240), text_size=10
    )
    cost_box = add_shape_box(
        slide, Inches(5.7), Inches(y_pos), Inches(1.8), Inches(0.5),
        cost, (255, 200, 200), text_size=10, bold=True
    )
    y_pos += 0.6

# Total investment
total_inv = add_shape_box(
    slide, Inches(0.5), Inches(y_pos), Inches(7), Inches(0.6),
    "TOTAL PHASE 1 INVESTMENT: $810K-1,140K",
    (255, 100, 100), text_size=12, bold=True
)

# Benefits section (RIGHT)
ben_label = add_shape_box(
    slide, Inches(8.5), Inches(1.2), Inches(7), Inches(0.5),
    "ANNUAL BENEFITS (Year 1+)", (230, 255, 230), text_size=14, bold=True
)

benefits = [
    ("Invoice Matching", "$280K-560K", "Month 4"),
    ("GL Anomaly Detector", "$192K", "Month 4"),
    ("Vendor Data Quality", "$200K-400K", "Month 6"),
    ("SoD Analyzer", "$40K", "Month 6"),
    ("e-Invoice Module", "Compliance", "Month 3"),
]

y_pos = 1.8
for item, value, timeline in benefits:
    ben_box = add_shape_box(
        slide, Inches(8.5), Inches(y_pos), Inches(4), Inches(0.5),
        item, (240, 255, 240), text_size=10
    )
    val_box = add_shape_box(
        slide, Inches(12.7), Inches(y_pos), Inches(1.5), Inches(0.5),
        value, (200, 255, 200), text_size=9, bold=True
    )
    time_box = add_shape_box(
        slide, Inches(14.3), Inches(y_pos), Inches(1.2), Inches(0.5),
        timeline, (220, 255, 220), text_size=8
    )
    y_pos += 0.6

# Total benefit
total_ben = add_shape_box(
    slide, Inches(8.5), Inches(y_pos), Inches(7), Inches(0.6),
    "TOTAL YEAR 1 BENEFIT: $722K-1.2M",
    (100, 255, 100), text_size=12, bold=True
)

# ROI Calculation (BOTTOM)
roi_label = add_shape_box(
    slide, Inches(0.5), Inches(5.5), Inches(15), Inches(0.5),
    "ROI CALCULATION", (255, 255, 200), text_size=14, bold=True
)

roi_metrics = [
    ("Year 1 Net Benefit", "-$278K to +$200K", "(Investment phase)"),
    ("Payback Period", "10-16 months", "(Strong business case)"),
    ("Year 2 Net Benefit", "+$722K-1.2M", "(Recurring profit)"),
    ("Year 3 Net Benefit", "+$722K-1.2M", "(Recurring profit)"),
    ("3-Year Total Net", "+$1.166M-2.1M", "(Excellent ROI)"),
    ("IRR", "45-65%", "(Very attractive)"),
]

y_pos = 6.1
for metric, value, note in roi_metrics:
    metric_box = add_shape_box(
        slide, Inches(0.5), Inches(y_pos), Inches(4.5), Inches(0.45),
        metric, (255, 255, 230), text_size=10, bold=True
    )
    val_box = add_shape_box(
        slide, Inches(5.2), Inches(y_pos), Inches(3), Inches(0.45),
        value, (255, 255, 180), text_size=10, bold=True
    )
    note_box = add_shape_box(
        slide, Inches(8.4), Inches(y_pos), Inches(3.1), Inches(0.45),
        note, (255, 255, 230), text_size=9
    )
    y_pos += 0.5

# Key takeaway
takeaway = add_shape_box(
    slide, Inches(12), Inches(6.1), Inches(3.5), Inches(2.3),
    "",
    (200, 240, 255), text_size=9
)
takeaway_text = takeaway.text_frame
takeaway_text.clear()
p = takeaway_text.paragraphs[0]
p.text = "KEY TAKEAWAYS"
p.font.bold = True
p.font.size = Pt(11)

points = [
    "✓ Strong ROI: 10-16 months",
    "✓ Low risk investment",
    "✓ Recurring annual benefit",
    "✓ Compliance deadline met",
    "✓ Scalable for multi-plant",
]
for point in points:
    p = takeaway_text.add_paragraph()
    p.text = point
    p.font.size = Pt(9)


# Save presentation
output_file = "/workspaces/cockpit/YTL_Cement_Native_PPT.pptx"
prs.save(output_file)
print(f"✓ Native PowerPoint presentation created: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"✓ All diagrams built with native PowerPoint shapes")
print(f"✓ Fully editable: colors, text, positions, sizes")
