#!/usr/bin/env python3
"""
Create YTL Cement PowerPoint Presentation from diagrams
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_diagram_slide(prs, title, image_path, notes=""):
    """Add a slide with diagram"""
    # Use blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Add image
    if os.path.exists(image_path):
        left = Inches(0.5)
        top = Inches(1.3)
        height = Inches(7)
        pic = slide.shapes.add_picture(image_path, left, top, height=height)

    # Add notes if provided
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

# Slide 1: Title
add_title_slide(
    prs,
    "YTL Cement - SAP S/4HANA Transformation",
    "Business Case & Technical Architecture\nGenerated with PlantUML"
)

# Slide 2: Capability Map AS-IS
add_diagram_slide(
    prs,
    "Business Capability Map - Current State (AS-IS)",
    "/workspaces/cockpit/YTL_Capability_Map_AS_IS.png",
    "Current state shows manual processes in RED with significant cost and risk issues"
)

# Slide 3: Capability Map TO-BE
add_diagram_slide(
    prs,
    "Business Capability Map - Target State (TO-BE)",
    "/workspaces/cockpit/YTL_Capability_Map_TO_BE.png",
    "Target state shows automated processes in GREEN with MVP modules providing intelligence"
)

# Slide 4: Application Landscape
add_diagram_slide(
    prs,
    "Application Landscape - AS-IS vs TO-BE",
    "/workspaces/cockpit/YTL_Application_Landscape.png",
    "Comparison of fragmented legacy systems vs integrated S/4HANA platform"
)

# Slide 5: Gap Analysis
add_diagram_slide(
    prs,
    "Gap Analysis - 5 Major Gaps & Solutions",
    "/workspaces/cockpit/YTL_Gap_Analysis.png",
    "5 critical gaps addressed by MVP modules with 10-16 month payback period"
)

# Slide 6: Data Flow
add_diagram_slide(
    prs,
    "Data Flow Architecture",
    "/workspaces/cockpit/YTL_Data_Flow.png",
    "End-to-end data flow from plant to marketing with MVP intelligence layer"
)

# Slide 7: Integration Architecture
add_diagram_slide(
    prs,
    "Integration Architecture - Middleware & APIs",
    "/workspaces/cockpit/YTL_Integration_Architecture.png",
    "Synchronous, asynchronous, and API integrations with SLA specifications"
)

# Slide 8: Technology Stack
add_diagram_slide(
    prs,
    "Technology Stack & Infrastructure",
    "/workspaces/cockpit/YTL_Technology_Stack.png",
    "Complete technology stack from application layer to infrastructure deployment"
)

# Slide 9: Implementation Roadmap
add_diagram_slide(
    prs,
    "Implementation Roadmap - 18 Week Phased Approach",
    "/workspaces/cockpit/YTL_Implementation_Roadmap.png",
    "Phase 1 (Weeks 1-12) Foundation & Quick Wins, Phase 2 (Weeks 13-18) Optimization"
)

# Slide 10: Organizational View
add_diagram_slide(
    prs,
    "Organization Structure & Stakeholder Matrix",
    "/workspaces/cockpit/YTL_Organizational_View.png",
    "Project organization with 46 plant users + 3 marketing users in Phase 1"
)

# Slide 11: ROI & Business Value
add_diagram_slide(
    prs,
    "ROI & Annual Business Value",
    "/workspaces/cockpit/YTL_ROI_Business_Value.png",
    "Total investment $810K-1,140K with $722K-1.2M annual benefit, 10-16 month payback"
)

# Save presentation
output_file = "/workspaces/cockpit/YTL_Cement_Presentation.pptx"
prs.save(output_file)
print(f"✓ PowerPoint presentation created: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
